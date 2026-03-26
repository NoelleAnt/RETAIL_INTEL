#!/usr/bin/env python
"""
Script to generate a PowerPoint presentation for the Retail ML project.
"""

import os
from pptx import Presentation
from pptx.util import Inches
from ml_models.model_trainer import run_notebook_analyses
import pandas as pd

def create_presentation():
    # Run analyses to generate plots and results
    print("Running analyses to gather data for presentation...")
    # Note: run_notebook_analyses doesn't return results, so we'll assume plots are saved
    run_notebook_analyses()

    # Create presentation
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Retail Data Mining and Machine Learning Project"
    subtitle.text = "Comprehensive Analysis of Retail Sales Data"

    # Data Preprocessing slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Data Preprocessing'
    tf = body_shape.text_frame
    tf.text = '• Loaded retail sales dataset\n• Handled missing values (median for numerical, mode for categorical)\n• Removed duplicates and corrected negative values\n• Converted dates and discretized features (Age_Group, Spending_Category, Price_Level)\n• Encoded categorical variables\n• Normalized numerical features using MinMaxScaler'

    # Clustering Analysis slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Clustering Analysis (K-Means)'
    tf = body_shape.text_frame
    tf.text = '• Determined optimal K using Elbow Method and Silhouette Score\n• Applied K-Means clustering\n• Analyzed cluster characteristics\n• Visualized clusters'

    # Add clustering plot
    plots_dir = os.path.join(os.path.dirname(__file__), 'ml_models', 'plots')
    img_path = os.path.join(plots_dir, 'clustering_elbow_silhouette.png')
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8))

    # Classification Analysis slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Classification Analysis'
    tf = body_shape.text_frame
    tf.text = '• Target: Spending Level (High/Low)\n• Models: Naïve Bayes and Decision Tree\n• Evaluated accuracy, classification reports, and feature importance\n• Generated confusion matrices'

    # Add confusion matrix plots
    img_path_nb = os.path.join(plots_dir, 'nb_confusion_matrix.png')
    if os.path.exists(img_path_nb):
        slide.shapes.add_picture(img_path_nb, Inches(1), Inches(2), width=Inches(4))

    img_path_dt = os.path.join(plots_dir, 'dt_confusion_matrix.png')
    if os.path.exists(img_path_dt):
        slide.shapes.add_picture(img_path_dt, Inches(5), Inches(2), width=Inches(4))

    # Association Mining slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Association Mining (FPGrowth)'
    tf = body_shape.text_frame
    tf.text = '• Prepared transaction data\n• Applied FPGrowth algorithm\n• Generated association rules with lift metric\n• Interpreted top rules for business insights'

    # Regression Analysis slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Regression Analysis'
    tf = body_shape.text_frame
    tf.text = '• Linear Regression on Total Amount\n• Evaluated R-squared and RMSE\n• Analyzed feature coefficients\n• Created actual vs predicted and residual plots'

    # Add regression plot
    img_path_reg = os.path.join(plots_dir, 'regression_plots.png')
    if os.path.exists(img_path_reg):
        slide.shapes.add_picture(img_path_reg, Inches(1), Inches(2), width=Inches(8))

    # Model Validation slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Model Validation and Summary'
    tf = body_shape.text_frame
    tf.text = '• Compared model performances\n• Validated results across analyses\n• Provided business implications\n• Ensured robustness and interpretability'

    # Conclusion slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Conclusion'
    tf = body_shape.text_frame
    tf.text = '• Successfully implemented comprehensive ML pipeline\n• From data preprocessing to model deployment\n• Web application for real-time predictions\n• Valuable insights for retail business decisions'

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), 'Retail_ML_Presentation.pptx')
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == '__main__':
    create_presentation()